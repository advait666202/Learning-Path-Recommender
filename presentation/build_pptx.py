"""
Generate the 10-slide presentation as a themed .pptx.

Styling follows DESIGN-linear.app.md: near-black canvas (#010102), light-gray
ink (#f7f8f8), a single lavender accent (#5e6ad2), generous negative-tracked
display headings. Each slide carries an eyebrow, a title, body content, a light
diagram built from native shapes, and full speaker notes.

Run:  uv run python presentation/build_pptx.py
Output: presentation/RDMU_Learning_Path_Recommender.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# Palette (Linear tokens).
# ---------------------------------------------------------------------------
CANVAS = RGBColor(0x01, 0x01, 0x02)
SURFACE_1 = RGBColor(0x0F, 0x10, 0x11)
SURFACE_2 = RGBColor(0x14, 0x15, 0x16)
HAIRLINE = RGBColor(0x23, 0x25, 0x2A)
PRIMARY = RGBColor(0x5E, 0x6A, 0xD2)
PRIMARY_HOVER = RGBColor(0x82, 0x8F, 0xFF)
INK = RGBColor(0xF7, 0xF8, 0xF8)
INK_MUTED = RGBColor(0xD0, 0xD6, 0xE0)
INK_SUBTLE = RGBColor(0x8A, 0x8F, 0x98)
SUCCESS = RGBColor(0x27, 0xA6, 0x44)

FONT = "Segoe UI"  # Inter/SF-Pro stand-in available on Windows

# 16:9 canvas.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(run, text, size, color, bold=False, italic=False, spacing=None):
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if spacing is not None:
        # Letter spacing (negative tracking) via XML.
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def bg(slide, color=CANVAS):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def card(slide, left, top, width, height, fill=SURFACE_1, line=HAIRLINE,
         radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    return box


def eyebrow(slide, text, left=Emu(685800), top=Emu(548640)):
    tf = add_textbox(slide, left, top, Emu(9000000), Emu(360000))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text.upper(), 12, PRIMARY_HOVER, bold=True, spacing=0.4)


def title(slide, text, left=Emu(685800), top=Emu(850000), size=40):
    tf = add_textbox(slide, left, top, Emu(10800000), Emu(1100000))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, size, INK, bold=True, spacing=-1.0)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullets(slide, left, top, width, height, items, size=15, gap=6):
    tf = add_textbox(slide, left, top, width, height)
    for i, (lead, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if lead:
            set_run(p.add_run(), f"{lead}  ", size, PRIMARY_HOVER, bold=True)
        set_run(p.add_run(), body, size, INK_MUTED)


def chip(slide, left, top, text, width=Emu(2750000), height=Emu(560000),
         fill=SURFACE_2, text_color=INK):
    box = card(slide, left, top, width, height, fill=fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(140000)
    tf.margin_right = Emu(120000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, 12.5, text_color, bold=True)
    return box


def arrow(slide, left, top, width=Emu(420000), height=Emu(360000)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.shadow.inherit = False
    _solid(a, PRIMARY)


CONTENT_LEFT = Emu(685800)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg(slide)
    return slide


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = new_slide(prs)
    # brand square
    sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(685800), Emu(2300000),
                            Emu(520000), Emu(520000))
    sq.shadow.inherit = False
    _solid(sq, PRIMARY)
    eyebrow(s, "RDMU · Reinforcement Learning", top=Emu(3050000))
    tf = add_textbox(s, CONTENT_LEFT, Emu(3350000), Emu(10800000), Emu(1700000))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Personalized Learning Path", 54, INK, bold=True, spacing=-2.0)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), "Recommender using Reinforcement Learning", 54, INK,
            bold=True, spacing=-2.0)
    sub = add_textbox(s, CONTENT_LEFT, Emu(5050000), Emu(10800000), Emu(700000))
    set_run(sub.paragraphs[0].add_run(),
            "MDP  ·  Q-Learning  ·  Exploration vs Exploitation  ·  Multi-Criteria Decision Making",
            18, INK_SUBTLE)
    notes(s,
          "This project answers a single, concrete question every adaptive tutor must "
          "solve: given where a learner is right now, which concept should they study "
          "next? We treat that decision as sequential — each choice changes the learner's "
          "state — and learn a policy with reinforcement learning, then let an instructor "
          "steer the final pick with transparent criteria weights. The whole system is "
          "reproducible (fixed seed 42) and ships as an 8-page dashboard.")


def slide_pipeline_strip(s, top):
    labels = ["State", "Q-values", "Top-k", "MCDM", "Next concept"]
    x = CONTENT_LEFT
    w = Emu(1900000)
    for i, lab in enumerate(labels):
        fill = PRIMARY if i == len(labels) - 1 else SURFACE_1
        tcol = RGBColor(0xFF, 0xFF, 0xFF) if i == len(labels) - 1 else INK
        chip(s, x, top, lab, width=w, fill=fill, text_color=tcol)
        x = Emu(x + w)
        if i < len(labels) - 1:
            arrow(s, Emu(x + 40000), Emu(top + 100000))
            x = Emu(x + 500000)


def slide_requirements(prs):
    s = new_slide(prs)
    eyebrow(s, "Requirement Understanding")
    title(s, "What the system must do")
    cols = [
        ("Adaptive", "Responds to the individual learner's state — prior knowledge, "
                     "study time, difficulty tolerance, interest."),
        ("Sequential", "Optimises a whole learning path, not a one-off pick — each "
                       "choice changes what comes next."),
        ("Explainable", "An instructor can see and adjust why a concept was chosen, "
                        "via interpretable criteria weights."),
    ]
    x = CONTENT_LEFT
    w = Emu(3500000)
    for name, body in cols:
        c = card(s, x, Emu(2150000), w, Emu(2400000))
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(220000)
        tf.margin_top = Emu(220000)
        tf.margin_right = Emu(200000)
        set_run(tf.paragraphs[0].add_run(), name, 22, INK, bold=True, spacing=-0.4)
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        set_run(p.add_run(), body, 14, INK_MUTED)
        x = Emu(x + w + Emu(330000))
    map_tf = add_textbox(s, CONTENT_LEFT, Emu(4850000), Emu(10800000), Emu(900000))
    set_run(map_tf.paragraphs[0].add_run(),
            "Satisfied by exactly four RDMU techniques: MDP, epsilon-greedy, tabular "
            "Q-Learning, and MCDM weighted scoring.", 15, INK_SUBTLE, italic=True)
    notes(s,
          "The requirement is a recommender that is adaptive, sequential and "
          "explainable. That rules out a static rules engine and a black-box "
          "classifier. It calls for a state-based decision model (MDP), a method that "
          "improves a policy from experience (Q-learning), a principled way to explore "
          "early (epsilon-greedy), and a transparent override layer (MCDM). We "
          "deliberately constrained ourselves to exactly these four RDMU techniques.")


def slide_business(prs):
    s = new_slide(prs)
    eyebrow(s, "Business Problem")
    title(s, "Why next-concept sequencing matters")
    bullets(s, CONTENT_LEFT, Emu(2100000), Emu(10800000), Emu(2600000), [
        ("Fixed curricula ignore the learner.",
         "Students differ in prior knowledge, study time, difficulty tolerance and interest."),
        ("Two failure modes.",
         "Concepts taught before their prerequisites cause frustration and drop-off; "
         "concepts too easy or too hard cause disengagement."),
        ("Measurable cost.",
         "Both lower course completion and weaken mastery."),
        ("Our goal.",
         "Personalise the ORDER of learning to each student's state — reaching the "
         "objective (Deep Learning) efficiently and at an appropriate difficulty, with "
         "the instructor in control via MCDM weights."),
    ], size=16, gap=12)
    notes(s,
          "On any learning platform, fixed linear curricula ignore that learners "
          "differ. Mis-sequencing causes two failure modes: prerequisites violated "
          "(frustration, drop-off) and difficulty mismatch (disengagement). The cost is "
          "measurable — lower completion and weaker mastery. Our system personalises the "
          "order of learning to each student's state, aiming to reach the objective "
          "efficiently and at an appropriate difficulty, with the instructor in control.")


def slide_flow(prs):
    s = new_slide(prs)
    eyebrow(s, "Application Flow Chart")
    title(s, "End-to-end application flow")
    stages = [
        ("1 · Data", "data_gen → students.csv\n2,200 students, seed 42"),
        ("2 · Calibrate + Train", "env calibrated from CSV;\nQ-learning trains online\n→ q_table.pkl"),
        ("3 · Serve", "app loads cached Q-table\n→ RL rank → top-k"),
        ("4 · Decide", "MCDM re-ranks shortlist\n→ final recommendation"),
    ]
    x = CONTENT_LEFT
    w = Emu(2550000)
    for i, (head, body) in enumerate(stages):
        c = card(s, x, Emu(2300000), w, Emu(1900000))
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(180000)
        tf.margin_top = Emu(170000)
        set_run(tf.paragraphs[0].add_run(), head, 15, PRIMARY_HOVER, bold=True)
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        set_run(p.add_run(), body, 13, INK_MUTED)
        x = Emu(x + w)
        if i < len(stages) - 1:
            arrow(s, Emu(x + 30000), Emu(3050000))
            x = Emu(x + 430000)
    note_tf = add_textbox(s, CONTENT_LEFT, Emu(4550000), Emu(10800000), Emu(800000))
    set_run(note_tf.paragraphs[0].add_run(),
            "Key separation: the agent trains against the simulated environment, never "
            "directly on the CSV. The dashboard never retrains on ordinary interaction.",
            14, INK_SUBTLE, italic=True)
    notes(s,
          "Data flows one direction with clean separation. A synthetic, internally-"
          "consistent dataset of 2,200 students calibrates the environment's transition "
          "probabilities and supplies profiles. The Q-learning agent trains online "
          "against the simulated environment — never directly on the CSV — and the "
          "learned Q-table is pickled. At inference the dashboard loads it (cached), "
          "ranks valid next concepts, hands the top-k to MCDM, and shows both rankings.")


def slide_concepts(prs):
    s = new_slide(prs)
    eyebrow(s, "Selected RDMU Concepts")
    title(s, "The four concepts, made concrete")
    cards = [
        ("A · Markov Decision Process",
         "State = (concept, mastery band, study-time band). Action = next concept "
         "(prerequisite-valid). Reward = weighted blend of 5 signals.  rl/environment.py"),
        ("B · Exploration vs Exploitation",
         "Epsilon-greedy selection; epsilon decays 1.0 → 0.1 floor. Explore early, "
         "exploit late.  rl/q_learning.py"),
        ("C · Tabular Q-Learning",
         "288×24 Q-table updated by temporal-difference learning over 5,000 simulated "
         "episodes.  rl/train.py"),
        ("D · Multi-Criteria Decision Making",
         "Weighted scoring over 5 instructor-tunable criteria re-ranks the RL shortlist; "
         "the winner is the final pick.  mcdm/scoring.py"),
    ]
    positions = [(CONTENT_LEFT, Emu(2100000)), (Emu(6450000), Emu(2100000)),
                 (CONTENT_LEFT, Emu(4150000)), (Emu(6450000), Emu(4150000))]
    w, h = Emu(5050000), Emu(1850000)
    for (left, top), (head, body) in zip(positions, cards):
        c = card(s, left, top, w, h)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(200000)
        tf.margin_top = Emu(160000)
        tf.margin_right = Emu(180000)
        set_run(tf.paragraphs[0].add_run(), head, 16, INK, bold=True, spacing=-0.3)
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        set_run(p.add_run(), body, 13, INK_MUTED)
    notes(s,
          "A — MDP: state is (current concept, mastery band, study-time band); action "
          "is the next concept restricted to valid candidates; reward blends mastery "
          "gain, difficulty suitability, preference match, time efficiency and progress. "
          "B — epsilon-greedy with epsilon decaying from 1.0 to a 0.1 floor. C — a "
          "288×24 Q-table updated by the TD rule across 5,000 episodes. D — weighted "
          "scoring over five tunable criteria re-ranks the RL shortlist; the winner is "
          "the final recommendation.")


def slide_dataset(prs):
    s = new_slide(prs)
    eyebrow(s, "Dataset Description")
    title(s, "Synthetic, consistent, reproducible")
    bullets(s, CONTENT_LEFT, Emu(2100000), Emu(6100000), Emu(3200000), [
        ("2,200 students, seed 42.", "Fully reproducible generation."),
        ("16 columns.", "Demographics, assessment scores, study habits, preferences, "
                        "and a prerequisite-valid completed→current→next concept triple."),
        ("Internally consistent.", "More study hours → higher mastery growth; quiz / "
                                   "assignment scores track current mastery."),
        ("Dual role.", "Calibrates the environment and drives analytics — the agent "
                       "never trains on it directly."),
    ], size=15, gap=11)
    # verified trend card
    c = card(s, Emu(7050000), Emu(2100000), Emu(4450000), Emu(2950000))
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(220000)
    tf.margin_top = Emu(200000)
    set_run(tf.paragraphs[0].add_run(), "Verified consistency", 15, PRIMARY_HOVER, bold=True)
    for lab, val, frac in [("Low study time", "0.34", 0.34),
                           ("Medium study time", "0.47", 0.47),
                           ("High study time", "0.70", 0.70)]:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        set_run(p.add_run(), f"{lab}", 13, INK_MUTED)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), "█" * int(frac * 22) + f"  {val}", 12, PRIMARY)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(10)
    set_run(p3.add_run(), "Mean mastery growth by study-time tercile", 11, INK_SUBTLE, italic=True)
    notes(s,
          "We generate 2,200 students with seed 42 so anyone can reproduce the exact "
          "dataset. The 16 columns span demographics, assessment scores, study habits, "
          "preferences and a prerequisite-valid concept triple from a topological order "
          "of the 24-concept DAG. The generative model is internally consistent: more "
          "weekly study hours yield higher mastery growth (verified terciles 0.34 / "
          "0.47 / 0.70), and scores track mastery. It calibrates the environment and "
          "drives analytics; the agent never trains on it directly.")


def slide_architecture(prs):
    s = new_slide(prs)
    eyebrow(s, "System Architecture")
    title(s, "Modular, with a single source of truth")
    layers = [
        ("utils/concepts.py", "Canonical 24-concept DAG — imported by graph, data, env"),
        ("utils/config.py", "Seed, band edges, reward & MCDM weights, hyperparameters, theme"),
        ("rl/  (environment · q_learning · train)", "MDP, epsilon-greedy agent, pickled Q-table"),
        ("mcdm/scoring.py", "Weighted-scoring re-ranker (isolated)"),
        ("visualizations/charts.py", "Pure Plotly figure builders"),
        ("app.py", "8-page Streamlit shell — cache_resource (model) + cache_data (data)"),
    ]
    y = Emu(2100000)
    for name, body in layers:
        c = card(s, CONTENT_LEFT, y, Emu(10800000), Emu(620000))
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(220000)
        p = tf.paragraphs[0]
        set_run(p.add_run(), f"{name}    ", 14, PRIMARY_HOVER, bold=True)
        set_run(p.add_run(), body, 13, INK_MUTED)
        y = Emu(y + Emu(720000))
    notes(s,
          "Every shared fact has one home. The 24-concept catalogue lives only in "
          "utils/concepts.py; graph, dataset and environment import it, so there's no "
          "drift. Configuration is centralised in utils/config.py. RL holds the "
          "environment, agent and trainer; MCDM is isolated; visualizations are pure "
          "functions returning Plotly figures; the app composes them with cache_resource "
          "for the model and cache_data for the data. Pre-training is a script, so the "
          "app loads instantly and never retrains on interaction.")


def slide_dashboard(prs):
    s = new_slide(prs)
    eyebrow(s, "Dashboard Screens")
    title(s, "Eight pages, one decision story")
    pages = ["1 · Project Overview", "2 · Student Profile", "3 · Concept Graph",
             "4 · Recommendation", "5 · RL Policy", "6 · Mastery Analytics",
             "7 · MCDM Settings", "8 · Performance Metrics"]
    x0, y0 = CONTENT_LEFT, Emu(2150000)
    w, h = Emu(2550000), Emu(1150000)
    gx, gy = Emu(330000), Emu(300000)
    for i, name in enumerate(pages):
        col = i % 4
        rowi = i // 4
        left = Emu(x0 + col * (w + gx))
        top = Emu(y0 + rowi * (h + gy))
        featured = (i == 3)  # Recommendation page
        c = card(s, left, top, w, h, fill=SURFACE_2 if featured else SURFACE_1,
                 line=PRIMARY if featured else HAIRLINE)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(180000)
        p = tf.paragraphs[0]
        set_run(p.add_run(), name, 14, INK if not featured else PRIMARY_HOVER,
                bold=featured)
    notes(s,
          "The dashboard walks from context to decision to evaluation. Overview frames "
          "the problem and the four concepts; Student Profile and the interactive "
          "Concept Graph give context; the Recommendation page is the heart — RL ranking "
          "and MCDM re-ranking side by side, plus the final pick; RL Policy exposes "
          "training internals (reward curve, epsilon decay, explore/exploit, policy "
          "stability, Q-value heatmap) and a retrain control; Mastery Analytics and MCDM "
          "Settings probe data and steer weights live; Performance Metrics closes with "
          "the KPI set. The visual language follows the Linear design system.")


def slide_results(prs):
    s = new_slide(prs)
    eyebrow(s, "Results & Metrics")
    title(s, "Does it actually learn?  Yes.")
    kpis = [
        ("Reward: early → converged", "1.54 → 8.49"),
        ("Path completion rate", "92%"),
        ("Policy stability", "98%"),
        ("Mastery rate (≥50% concepts)", "31%"),
        ("Reco. accuracy (RL≡MCDM proxy)", "51%"),
        ("State space", "288"),
    ]
    x0, y0 = CONTENT_LEFT, Emu(2150000)
    w, h = Emu(3500000), Emu(1300000)
    gx, gy = Emu(330000), Emu(300000)
    for i, (lab, val) in enumerate(kpis):
        col = i % 3
        rowi = i // 3
        left = Emu(x0 + col * (w + gx))
        top = Emu(y0 + rowi * (h + gy))
        c = card(s, left, top, w, h)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(200000)
        tf.margin_top = Emu(170000)
        set_run(tf.paragraphs[0].add_run(), lab.upper(), 11, INK_SUBTLE, bold=True, spacing=0.3)
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), val, 30, INK, bold=True, spacing=-1.0)
    notes(s,
          "Yes — measured with precisely-defined metrics, not a vague accuracy. Mean "
          "episodic reward rose from 1.54 (early, exploratory) to 8.49 (converged) — a "
          "clear learning curve produced by a deliberately tight 22-step budget that "
          "forces efficiency. The greedy policy reaches the objective on ~92% of paths, "
          "and ~98% of states have a stable greedy action across the last checkpoints. "
          "Recommendation Accuracy — our stated RL-vs-MCDM agreement proxy — is ~51%, the "
          "point being that MCDM genuinely re-ranks about half the time, so the "
          "instructor's weights matter. Mastery rate (students past 50% of the "
          "curriculum) is about 31%.")


def slide_conclusion(prs):
    s = new_slide(prs)
    eyebrow(s, "Conclusion & Future Scope")
    title(s, "Conclusion and what's next")
    left = card(s, CONTENT_LEFT, Emu(2150000), Emu(5250000), Emu(3100000))
    tf = left.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(220000)
    tf.margin_top = Emu(200000)
    set_run(tf.paragraphs[0].add_run(), "Delivered", 17, SUCCESS, bold=True)
    for b in ["Reproducible, runnable end-to-end prototype.",
              "Exactly four RDMU techniques, cleanly separated.",
              "RL (what tends to work) + MCDM (human judgement), both visible.",
              "8-page dashboard with precise metrics."]:
        p = tf.add_paragraph()
        p.space_before = Pt(9)
        set_run(p.add_run(), "•  ", 14, SUCCESS)
        set_run(p.add_run(), b, 14, INK_MUTED)
    right = card(s, Emu(6250000), Emu(2150000), Emu(5250000), Emu(3100000))
    tf2 = right.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Emu(220000)
    tf2.margin_top = Emu(200000)
    set_run(tf2.paragraphs[0].add_run(), "Future scope", 17, PRIMARY_HOVER, bold=True)
    for b in ["Function approximation (Deep RL) to scale beyond discrete bands.",
              "Calibrate on real learner telemetry instead of synthetic data.",
              "A/B evaluation against a fixed-curriculum baseline.",
              "Richer MCDM criteria (spaced repetition, cohort outcomes)."]:
        p = tf2.add_paragraph()
        p.space_before = Pt(9)
        set_run(p.add_run(), "•  ", 14, PRIMARY)
        set_run(p.add_run(), b, 14, INK_MUTED)
    notes(s,
          "We delivered a reproducible, examination-grade prototype that uses exactly "
          "the four RDMU techniques to make an explainable, adaptive next-concept "
          "recommendation, end to end and fully runnable. The design keeps RL and human "
          "judgement (MCDM weights) cleanly separated and visible. Future scope: replace "
          "the tabular Q-table with function approximation to scale; calibrate on real "
          "learner telemetry; add A/B evaluation against a fixed-curriculum baseline; and "
          "extend MCDM to richer criteria. None of these change the core RDMU framing — "
          "they deepen each layer.")


def build() -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    # add pipeline strip to the title slide bottom
    slide_pipeline_strip(prs.slides[0], Emu(5750000))

    slide_requirements(prs)
    slide_business(prs)
    slide_flow(prs)
    slide_concepts(prs)
    slide_dataset(prs)
    slide_architecture(prs)
    slide_dashboard(prs)
    slide_results(prs)
    slide_conclusion(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "RDMU_Learning_Path_Recommender.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"Saved 10 slides to {path}")
